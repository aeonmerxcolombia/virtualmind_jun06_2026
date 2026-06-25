# app/routes/documentos_office_router.py - Router para documentos Office con OnlyOffice

from fastapi import (
    APIRouter,
    Depends,
    HTTPException,
    status,
    UploadFile,
    File,
    Form,
    Request,
    Query,
)
import requests
import jwt
from sqlalchemy.orm import Session
from sqlalchemy import text
from typing import List, Optional
from datetime import datetime
import os
from io import BytesIO

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import re
import io

from pydantic import BaseModel

from app.database.db import SessionLocal
from app.auth.jwt_handler import verify_token
from app.services.log_service import registrar_log, crear_notificacion
from app.services.email_service import send_email
from app.models.user import User
from app.models.project import Project
from app.models.archivo import Archivo

# Importamos los nuevos Modelos y Esquemas
from app.models.documento_office import (
    DocumentoOficina,
    DocumentoOficinaCompartido,
    DocumentoVersion,
    DocumentoRevision,
    EstadoVersion,
)
from app.models.biblioteca import DocumentoBiblioteca, AccesoBiblioteca
from app.models.proyecto_participante import ProyectoParticipante
from app.routes.biblioteca_router import generar_nota_bibliografica
from app.schemas.documento_office_schema import DocumentoOfficeOut

router = APIRouter(prefix="/documentos-office", tags=["Documentos Office"])

DOCUMENTOS_DIR = "/home/ubuntu/backend/documentos/archivos"
os.makedirs(DOCUMENTOS_DIR, exist_ok=True)
ONLYOFFICE_SECRET = "pQgLEoIAyJEVcFrHkw5OoL6i6htHMs5Q"


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


# ==================== GENERADORES NATIVOS ====================
def _generate_docx_template(nombre: str) -> bytes:
    doc = Document()
    doc.add_heading(nombre, 0)
    doc.add_paragraph("Documento creado con VirtualMind")
    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _generate_xlsx_template(nombre: str) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Hoja 1"
    ws["A1"] = nombre
    ws["A2"] = "Documento creado con VirtualMind"
    ws.column_dimensions["A"].width = 30
    buffer = BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def _generate_pptx_template(nombre: str) -> bytes:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = nombre
    slide.placeholders[1].text = "Presentación creada con VirtualMind"
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ==================== ENDPOINTS ====================


@router.get("/", response_model=List[DocumentoOfficeOut])
def list_documentos(
    project_id: Optional[int] = Query(None),
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    query = db.query(DocumentoOficina)

    user_email = token_data.get("sub")
    user_uid = token_data.get("user_id") or user_email
    roles = token_data.get("roles", [])

    if "superadmin" not in roles:
        query = query.filter(DocumentoOficina.usuario_id == user_uid)

    if project_id:
        query = query.filter(DocumentoOficina.project_id == project_id)

    docs = query.all()
    return docs


@router.post(
    "/", response_model=DocumentoOfficeOut, status_code=status.HTTP_201_CREATED
)
def crear_documento(
    nombre: str = Form(...),
    tipo: str = Form(...),
    project_id: Optional[int] = Form(None),  # NUEVO
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    if tipo not in ["docx", "xlsx", "pptx"]:
        raise HTTPException(
            status_code=400, detail="Tipo no soportado. Usa: docx, xlsx, pptx"
        )

    user_uid = str(token_data.get("user_id") or token_data.get("sub") or "anon_user")

    # Crear en DB primero para obtener ID
    now = datetime.now()
    doc = DocumentoOficina(
        project_id=project_id,
        nombre=nombre,
        tipo=tipo,
        filename="temp",
        ruta="",
        usuario_id=user_uid,
        creado=now,
        actualizado=now,
    )
    db.add(doc)
    db.flush()  # Genera el ID

    filename = f"{nombre}_{doc.id}.{tipo}"
    filepath = os.path.join(DOCUMENTOS_DIR, filename)

    if tipo == "docx":
        content = _generate_docx_template(nombre)
    elif tipo == "xlsx":
        content = _generate_xlsx_template(nombre)
    elif tipo == "pptx":
        content = _generate_pptx_template(nombre)

    with open(filepath, "wb") as f:
        f.write(content)

    doc.filename = filename
    doc.ruta = f"/documentos-office/archivo/{filename}"
    doc.url_editar = f"/r/superadmin/documentos/editor.html?id={doc.id}"

    db.commit()
    db.refresh(doc)

    return doc


@router.get("/documento/{doc_id}", response_model=DocumentoOfficeOut)
async def get_documento(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="No encontrado")

    # Notificar al propietario si quien abre no es el dueño
    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_nombre = token_data.get("nombre", "Usuario")

    if doc.usuario_id != user_uid:
        # Verificar que esté compartido con este usuario
        compartido = (
            db.query(DocumentoOficinaCompartido)
            .filter(
                DocumentoOficinaCompartido.documento_id == doc_id,
                DocumentoOficinaCompartido.email_destino == token_data.get("sub"),
            )
            .first()
        )

        if compartido:
            # Notificar al dueño que abrieron su documento
            try:
                from datetime import datetime

                owner = (
                    db.query(User)
                    .filter(
                        (User.uid == doc.usuario_id) | (User.email == doc.usuario_id)
                    )
                    .first()
                )

                if owner:
                    crear_notificacion(
                        db=db,
                        usuario_id=owner.uid,
                        tipo_evento="documento_abierto",
                        descripcion=f"{user_nombre} abrió el documento '{doc.nombre}' que compartiste",
                        link=f"/r/autor/documentos/index.html?doc={doc_id}",
                    )

                    # Enviar correo al dueño
                    try:
                        asunto = (
                            f"{user_nombre} abrió tu documento compartido: {doc.nombre}"
                        )
                        cuerpo = f"""Hola,

{user_nombre} ha abierto el documento "{doc.nombre}" que compartiste con él/ella.

Fecha: {datetime.now().strftime("%d/%m/%Y %H:%M")}

Saludos,
VirtualMind"""
                        await send_email(owner.email, asunto, cuerpo)
                    except Exception as e:
                        print(f"Error enviando correo al dueño: {e}")
            except Exception as e:
                print(f"Error notificando apertura: {e}")

            # Marcar como "visto" en la revisión correspondiente
            try:
                roles = token_data.get("roles", [])
                revision = (
                    db.query(DocumentoRevision)
                    .filter(
                        DocumentoRevision.documento_id == doc_id,
                        DocumentoRevision.estado.in_(
                            ["pendiente", "aprobado", "enviado_cliente"]
                        ),
                    )
                    .order_by(DocumentoRevision.id.desc())
                    .first()
                )

                if revision:
                    now = datetime.now()
                    if "coordinador" in roles and not revision.coordinador_visto:
                        revision.coordinador_visto = now
                    elif "cliente" in roles and not revision.cliente_visto:
                        revision.cliente_visto = now
                    db.commit()
            except Exception as e:
                print(f"Error marcando visto: {e}")

    return doc


# ==================== VERSIONADO ====================

VERSIONES_DIR = "/home/ubuntu/backend/documentos/versiones"
os.makedirs(VERSIONES_DIR, exist_ok=True)


@router.post("/version/{doc_id}/nueva")
async def crear_nueva_version(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    if doc.estado_version == "cerrada":
        raise HTTPException(
            status_code=400,
            detail="No se puede crear una nueva versión porque la versión actual está cerrada. La versión cerrada es definitiva.",
        )

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_email = token_data.get("sub", "")

    # Archivar versión actual
    versiones_dir_doc = os.path.join(VERSIONES_DIR, str(doc_id))
    os.makedirs(versiones_dir_doc, exist_ok=True)

    current_file = os.path.join(DOCUMENTOS_DIR, doc.filename)
    if os.path.exists(current_file):
        version_actual = doc.version_actual or "1.0"
        archivo_version = f"v{version_actual}_{doc.filename}"
        import shutil

        shutil.copy2(current_file, os.path.join(versiones_dir_doc, archivo_version))

        # Registrar versión en DB
        ver = DocumentoVersion(
            documento_id=doc_id,
            version=version_actual,
            estado=doc.estado_version or "en_progreso",
            archivo_filename=archivo_version,
            creado_por=user_uid,
            comentarios=f"Versión {version_actual} archivada por {user_email}",
        )
        db.add(ver)

    # Calcular nueva versión
    partes = (doc.version_actual or "1.0").split(".")
    try:
        nueva_mayor = int(partes[0]) + 1
    except:
        nueva_mayor = 2
    nueva_version = f"{nueva_mayor}.0"
    doc.version_actual = nueva_version
    doc.estado_version = "en_progreso"

    db.commit()
    db.refresh(doc)

    return {
        "message": f"Nueva versión {nueva_version} creada",
        "version_actual": doc.version_actual,
        "estado_version": doc.estado_version,
        "documento_id": doc.id,
    }


class CerrarVersionRequest(BaseModel):
    roles_acceso: List[str] = []
    usuarios_acceso: List[str] = []


@router.post("/version/{doc_id}/cerrar")
async def cerrar_version(
    doc_id: int,
    body: Optional[CerrarVersionRequest] = None,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    roles = token_data.get("roles", [])

    # Solo autor y cliente pueden cerrar versiones
    roles_permitidos = {"autor", "cliente", "superadmin"}
    if not roles_permitidos.intersection(set(roles)) and "superadmin" not in roles:
        raise HTTPException(
            status_code=403,
            detail="Solo Autor, Cliente o Superadmin pueden cerrar versiones",
        )

    if doc.estado_version == "cerrada":
        raise HTTPException(status_code=400, detail="Esta versión ya está cerrada")

    # Archivar la versión actual antes de cerrar
    versiones_dir_doc = os.path.join(VERSIONES_DIR, str(doc_id))
    os.makedirs(versiones_dir_doc, exist_ok=True)

    current_file = os.path.join(DOCUMENTOS_DIR, doc.filename)
    if os.path.exists(current_file):
        version_actual = doc.version_actual or "1.0"
        archivo_version = f"v{version_actual}_final_{doc.filename}"
        import shutil

        shutil.copy2(current_file, os.path.join(versiones_dir_doc, archivo_version))

        from datetime import datetime

        ver = DocumentoVersion(
            documento_id=doc_id,
            version=version_actual,
            estado="cerrada",
            archivo_filename=archivo_version,
            creado_por=user_uid,
            comentarios=f"Versión {version_actual} cerrada/aprobada",
            fecha_cierre=datetime.now(),
        )
        db.add(ver)

    doc.estado_version = "cerrada"

    # Enviar automáticamente a la Biblioteca de Documentos
    BIBLIOTECA_DIR = "/home/ubuntu/backend/documentos/biblioteca"
    os.makedirs(BIBLIOTECA_DIR, exist_ok=True)
    biblioteca_filename = f"bib_{doc.id}_v{version_actual}_{doc.filename}"
    biblioteca_filepath = os.path.join(BIBLIOTECA_DIR, biblioteca_filename)

    if os.path.exists(current_file):
        import shutil

        shutil.copy2(current_file, biblioteca_filepath)

    # Obtener nombre del proyecto si existe
    proyecto_nombre = None
    if doc.project_id:
        try:
            proyecto = db.query(Project).filter(Project.id == doc.project_id).first()
            if proyecto:
                proyecto_nombre = proyecto.nombre
        except:
            pass

    # Obtener nombre del usuario
    user_nombre = token_data.get("nombre", "Usuario")

    # Generar nota bibliográfica
    from datetime import datetime

    year = datetime.now().year
    nota = generar_nota_bibliografica(
        nombre=doc.nombre,
        autor=user_nombre,
        proyecto=proyecto_nombre or "Sin proyecto",
        tipo=doc.tipo,
        year=year,
    )

    bib_entry = DocumentoBiblioteca(
        documento_id=doc.id,
        project_id=doc.project_id,
        nombre=doc.nombre,
        tipo=doc.tipo,
        version=version_actual,
        filename=biblioteca_filename,
        ruta_archivo=biblioteca_filepath,
        usuario_id=user_uid,
        usuario_nombre=user_nombre,
        proyecto_nombre=proyecto_nombre,
        nota_bibliografica=nota,
        descripcion=f"Documento '{doc.nombre}' versión {version_actual} cerrado y archivado en biblioteca",
        etiquetas=f"{doc.tipo}, v{version_actual}",
    )
    db.add(bib_entry)
    db.flush()

    # Otorgar acceso en biblioteca
    try:
        # Siempre dar acceso al creador
        acceso_creador = AccesoBiblioteca(
            documento_id=bib_entry.id,
            usuario_id=user_uid,
            permiso="lectura",
        )
        db.add(acceso_creador)

        roles_seleccionados = body.roles_acceso if body and body.roles_acceso else []
        usuarios_seleccionados = (
            body.usuarios_acceso if body and body.usuarios_acceso else []
        )

        if roles_seleccionados or usuarios_seleccionados:
            # Usar selección explícita del usuario
            for rol in roles_seleccionados:
                db.add(
                    AccesoBiblioteca(
                        documento_id=bib_entry.id,
                        rol=rol,
                        permiso="lectura",
                    )
                )
            for uid in usuarios_seleccionados:
                if uid and uid != user_uid:
                    existe = (
                        db.query(AccesoBiblioteca)
                        .filter(
                            AccesoBiblioteca.documento_id == bib_entry.id,
                            AccesoBiblioteca.usuario_id == uid,
                        )
                        .first()
                    )
                    if not existe:
                        db.add(
                            AccesoBiblioteca(
                                documento_id=bib_entry.id,
                                usuario_id=uid,
                                permiso="lectura",
                            )
                        )
        elif doc.project_id:
            # Fallback: participantes del proyecto
            participantes = (
                db.query(ProyectoParticipante)
                .filter(ProyectoParticipante.project_id == doc.project_id)
                .all()
            )
            for p in participantes:
                if p.user_uid != user_uid:
                    existe = (
                        db.query(AccesoBiblioteca)
                        .filter(
                            AccesoBiblioteca.documento_id == bib_entry.id,
                            AccesoBiblioteca.usuario_id == p.user_uid,
                        )
                        .first()
                    )
                    if not existe:
                        db.add(
                            AccesoBiblioteca(
                                documento_id=bib_entry.id,
                                usuario_id=p.user_uid,
                                permiso="lectura",
                            )
                        )
    except Exception as e:
        print(f"Error otorgando accesos biblioteca: {e}")

    db.commit()
    db.refresh(doc)

    return {
        "message": f"Versión {doc.version_actual} cerrada exitosamente y enviada a la biblioteca",
        "version_actual": doc.version_actual,
        "estado_version": doc.estado_version,
        "documento_id": doc.id,
        "biblioteca_id": bib_entry.id,
        "nota_bibliografica": nota,
        "biblioteca_nombre": doc.nombre,
    }


@router.get("/version/{doc_id}")
async def listar_versiones(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    versiones = (
        db.query(DocumentoVersion)
        .filter(DocumentoVersion.documento_id == doc_id)
        .order_by(DocumentoVersion.id.desc())
        .all()
    )

    return [
        {
            "id": v.id,
            "version": v.version,
            "estado": v.estado,
            "comentarios": v.comentarios,
            "fecha_creacion": v.fecha_creacion.isoformat()
            if v.fecha_creacion
            else None,
            "fecha_cierre": v.fecha_cierre.isoformat() if v.fecha_cierre else None,
            "creado_por": v.creado_por,
            "es_actual": v.version == doc.version_actual,
        }
        for v in versiones
    ]


@router.get("/version/archivo/{version_id}")
async def descargar_version(
    version_id: int,
    token: Optional[str] = Query(None),
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    ver = db.query(DocumentoVersion).filter(DocumentoVersion.id == version_id).first()
    if not ver:
        raise HTTPException(status_code=404, detail="Versión no encontrada")

    filepath = os.path.join(VERSIONES_DIR, str(ver.documento_id), ver.archivo_filename)
    if not os.path.exists(filepath):
        raise HTTPException(
            status_code=404, detail="Archivo de versión no encontrado en disco"
        )

    from fastapi.responses import FileResponse

    return FileResponse(filepath, filename=ver.archivo_filename)


@router.delete("/eliminar/{doc_id}", status_code=status.HTTP_204_NO_CONTENT)
def eliminar_documento(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="No encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    roles = token_data.get("roles", [])

    if "superadmin" not in roles and doc.usuario_id != user_uid:
        raise HTTPException(status_code=403, detail="Sin permiso")

    filepath = os.path.join(DOCUMENTOS_DIR, doc.filename)
    if os.path.exists(filepath):
        try:
            os.remove(filepath)
        except:
            pass

    db.delete(doc)
    db.commit()
    return None


# ==================== ONLYOFFICE & UPLOADS ====================


@router.get("/archivo/{filename}")
async def get_archivo(filename: str):
    filepath = os.path.join(DOCUMENTOS_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    from fastapi.responses import FileResponse

    ext = filename.split(".")[-1].lower()
    media_types = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    response = FileResponse(
        filepath,
        media_type=media_types.get(ext, "application/octet-stream"),
        filename=filename,
    )
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["X-Frame-Options"] = "ALLOWALL"
    return response


@router.post("/callback/{filename}")
async def onlyoffice_callback(
    filename: str, request: Request, db: Session = Depends(get_db)
):
    try:
        body = await request.json()
        status = body.get("status")

        if status == 2 or status == 6:
            download_url = body.get("url")
            if download_url:
                response = requests.get(download_url)
                if response.status_code == 200:
                    filepath = os.path.join(DOCUMENTOS_DIR, filename)
                    with open(filepath, "wb") as f:
                        f.write(response.content)

                    # Update MySQL
                    doc = (
                        db.query(DocumentoOficina)
                        .filter(DocumentoOficina.filename == filename)
                        .first()
                    )
                    if doc:
                        doc.actualizado = datetime.now()
                        db.commit()

        return {"error": 0}
    except Exception as e:
        return {"error": 1, "message": str(e)}


class ConfigArchivoRequest(BaseModel):
    file_id: int
    file_url: str
    file_name: str
    file_type: str


@router.post("/config-archivo")
async def config_archivo_from_drive(
    req: ConfigArchivoRequest,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    arch = db.query(Archivo).filter(Archivo.id == req.file_id).first()
    if not arch:
        raise HTTPException(status_code=404, detail="Archivo no encontrado en drive")

    ext = req.file_type.lower().replace(".", "")
    if ext not in ("docx", "xlsx", "pptx"):
        raise HTTPException(
            status_code=400, detail="Tipo de archivo no soportado para OnlyOffice"
        )

    # Copiar archivo a DOCUMENTOS_DIR
    src_filename = arch.url.rsplit("/", 1)[-1]
    src_path = os.path.join("/home/ubuntu/backend/static/uploads", src_filename)
    if not os.path.exists(src_path):
        raise HTTPException(
            status_code=404, detail="Archivo físico no encontrado en servidor"
        )

    user_uid = str(token_data.get("user_id") or token_data.get("sub") or "anon_user")
    doc = DocumentoOficina(
        nombre=os.path.splitext(req.file_name)[0],
        tipo=ext,
        filename="temp",
        ruta="",
        usuario_id=user_uid,
    )
    db.add(doc)
    db.flush()

    dest_filename = f"{os.path.splitext(req.file_name)[0]}_{doc.id}.{ext}"
    dest_path = os.path.join(DOCUMENTOS_DIR, dest_filename)
    import shutil

    shutil.copy2(src_path, dest_path)

    doc.filename = dest_filename
    doc.ruta = dest_path
    db.commit()

    # Generar config de OnlyOffice (misma lógica que config/{doc_id})
    base_url = "https://gestordecursos.pegui.edu.co:8000"
    document_type_map = {"docx": "word", "xlsx": "cell", "pptx": "slide"}
    key_time = int(doc.actualizado.timestamp())

    user_name = str(token_data.get("nombre") or token_data.get("sub", "").split("@")[0])

    config = {
        "document": {
            "fileType": doc.tipo,
            "key": f"{doc.id}_{key_time}",
            "title": f"{doc.nombre}.{doc.tipo}  (v{doc.version_actual})",
            "url": f"{base_url}/documentos-office/archivo/{doc.filename}",
            "permissions": {"download": True, "edit": True, "print": True},
        },
        "documentType": document_type_map.get(doc.tipo, "word"),
        "editorConfig": {
            "callbackUrl": f"{base_url}/documentos-office/callback/{doc.filename}",
            "lang": "es",
            "mode": "edit",
            "user": {"id": user_uid, "name": user_name},
            "customization": {
                "forcesave": True,
                "autosave": True,
                "about": False,
                "feedback": False,
            },
        },
    }
    token = jwt.encode(config, ONLYOFFICE_SECRET, algorithm="HS256")
    config["token"] = token if isinstance(token, str) else token.decode("utf-8")
    return {"config": config}


@router.get("/config/{doc_id}")
async def get_onlyoffice_config(
    doc_id: int,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="No encontrado")

    base_url = "https://gestordecursos.pegui.edu.co:8000"
    document_type_map = {"docx": "word", "xlsx": "cell", "pptx": "slide"}
    key_time = int(doc.actualizado.timestamp())

    user_uid = str(token_data.get("user_id") or token_data.get("sub") or "anon_user")
    user_name = str(token_data.get("nombre") or token_data.get("sub", "").split("@")[0])

    # Si la versión está cerrada, solo permitir lectura (excepto superadmin)
    roles = token_data.get("roles", [])
    puede_editar = doc.estado_version != "cerrada" or "superadmin" in roles

    # Marcar como visto si es coordinador o cliente
    if doc.usuario_id != user_uid:
        try:
            revision = (
                db.query(DocumentoRevision)
                .filter(
                    DocumentoRevision.documento_id == doc_id,
                    DocumentoRevision.estado.in_(
                        ["pendiente", "aprobado", "enviado_cliente"]
                    ),
                )
                .order_by(DocumentoRevision.id.desc())
                .first()
            )

            if revision:
                from datetime import datetime

                now = datetime.now()
                if "coordinador" in roles and not revision.coordinador_visto:
                    revision.coordinador_visto = now
                elif "cliente" in roles and not revision.cliente_visto:
                    revision.cliente_visto = now
                db.commit()
        except Exception as e:
            print(f"Error marcando visto en config: {e}")

    config = {
        "document": {
            "fileType": doc.tipo,
            "key": f"{doc.id}_{key_time}",
            "title": f"{doc.nombre}.{doc.tipo}  (v{doc.version_actual})",
            "url": f"{base_url}/documentos-office/archivo/{doc.filename}",
            "permissions": {"download": True, "edit": puede_editar, "print": True},
        },
        "documentType": document_type_map.get(doc.tipo, "word"),
        "editorConfig": {
            "callbackUrl": f"{base_url}/documentos-office/callback/{doc.filename}",
            "lang": "es",
            "mode": "edit" if puede_editar else "view",
            "user": {"id": user_uid, "name": user_name},
            "customization": {
                "forcesave": True,
                "autosave": True,
                "about": False,
                "feedback": False,
            },
        },
    }
    token = jwt.encode(config, ONLYOFFICE_SECRET, algorithm="HS256")
    config["token"] = token if isinstance(token, str) else token.decode("utf-8")
    return {"config": config}


@router.post("/upload")
@router.post("/subir-archivo")
async def upload_documento_v2(
    archivo: UploadFile = File(...),
    project_id: Optional[int] = Form(None),  # NUEVO
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    ext = archivo.filename.split(".")[-1].lower() if "." in archivo.filename else "docx"
    user_uid = str(token_data.get("user_id") or token_data.get("sub") or "anon_user")

    content = await archivo.read()
    nombre = os.path.splitext(archivo.filename)[0]

    doc = DocumentoOficina(
        project_id=project_id,
        nombre=nombre,
        tipo=ext,
        filename="temp",
        ruta="",
        usuario_id=user_uid,
    )
    db.add(doc)
    db.flush()

    filename = f"{archivo.filename}_{doc.id}"
    filepath = os.path.join(DOCUMENTOS_DIR, filename)

    with open(filepath, "wb") as f:
        f.write(content)

    doc.filename = filename
    doc.ruta = filepath
    db.commit()

    return {
        "message": "Archivo subido exitosamente",
        "id": doc.id,
        "nombre": nombre,
        "tipo": ext,
    }


@router.get("/mis-compartidos-usuario", response_model=List[DocumentoOfficeOut])
def mis_documentos_compartidos(
    db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    user_email = token_data.get("sub")
    if not user_email:
        raise HTTPException(status_code=401, detail="No autorizado")

    compartidos = (
        db.query(DocumentoOficina)
        .join(
            DocumentoOficinaCompartido,
            DocumentoOficina.id == DocumentoOficinaCompartido.documento_id,
        )
        .filter(DocumentoOficinaCompartido.email_destino == user_email)
        .all()
    )

    return compartidos


@router.post("/compartir/{doc_id}")
async def compartir_documento(
    doc_id: int,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    try:
        body = await request.json()
    except:
        raise HTTPException(status_code=422, detail="JSON inválido")

    email_destino = body.get("email_destino")
    permiso = body.get("permiso", "lectura")

    if not email_destino:
        raise HTTPException(status_code=422, detail="email_destino es requerido")

    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_email = token_data.get("sub", "")
    user_nombre = token_data.get("nombre", "Usuario")
    roles = token_data.get("roles", [])

    if "superadmin" not in roles and doc.usuario_id != user_uid:
        raise HTTPException(status_code=403, detail="Sin permiso")

    existente = (
        db.query(DocumentoOficinaCompartido)
        .filter(
            DocumentoOficinaCompartido.documento_id == doc_id,
            DocumentoOficinaCompartido.email_destino == email_destino,
        )
        .first()
    )

    if existente:
        raise HTTPException(
            status_code=400, detail="Ya está compartido con este usuario"
        )

    compartido = DocumentoOficinaCompartido(
        documento_id=doc_id, email_destino=email_destino, permiso=permiso
    )
    db.add(compartido)
    db.commit()

    # Obtener usuario destinatario
    destinatario = db.query(User).filter(User.email == email_destino).first()

    # Enviar correo al destinatario
    try:
        from datetime import datetime

        enlace = f"https://gestordecursos.pegui.edu.co/r/autor/documentos/index.html?doc={doc_id}"
        asunto = f"{user_nombre} te ha compartido un documento: {doc.nombre}"
        cuerpo = f"""Hola,

{user_nombre} te ha compartido el documento "{doc.nombre}".

Puedes acceder al documento aquí:
{enlace}

Fecha: {datetime.now().strftime("%d/%m/%Y %H:%M")}

Saludos,
VirtualMind"""
        import asyncio

        asyncio.create_task(send_email(email_destino, asunto, cuerpo))
    except Exception as e:
        print(f"Error enviando correo: {e}")

    # Crear notificación para el destinatario
    if destinatario:
        try:
            crear_notificacion(
                db=db,
                usuario_id=destinatario.uid,
                tipo_evento="documento_compartido",
                descripcion=f"{user_nombre} te compartió el documento '{doc.nombre}'",
                link=f"/r/autor/documentos/index.html?doc={doc_id}",
            )
        except Exception as e:
            print(f"Error creando notificación: {e}")

    return {"message": "Documento compartido exitosamente"}


# ==================== REVISIÓN (COORDINADOR) ====================


@router.post("/revision/{doc_id}/enviar")
async def enviar_a_revision(
    doc_id: int,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    """Autor envía documento al Coordinador para revisión"""
    try:
        body = await request.json()
    except:
        body = {}

    coordinador_email = body.get("coordinador_email", "")
    if not coordinador_email:
        raise HTTPException(
            status_code=422, detail="email del coordinador es requerido"
        )

    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_email = token_data.get("sub", "")
    user_nombre = token_data.get("nombre", "Usuario")

    if doc.estado_version == "cerrada":
        raise HTTPException(
            status_code=400, detail="No se puede enviar una versión cerrada a revisión"
        )

    # Buscar coordinador por email y validar rol
    coordinador = db.query(User).filter(User.email == coordinador_email).first()
    if not coordinador:
        raise HTTPException(status_code=404, detail="Coordinador no encontrado")
    es_coordinador = any(r.name == "coordinador" for r in (coordinador.roles or []))
    es_superadmin = any(r.name == "superadmin" for r in (coordinador.roles or []))
    if not es_coordinador and not es_superadmin:
        raise HTTPException(
            status_code=403,
            detail="El usuario seleccionado no tiene rol de Coordinador",
        )

    # Registrar la solicitud de revisión
    revision = DocumentoRevision(
        documento_id=doc_id,
        version=doc.version_actual or "1.0",
        autor_id=user_uid,
        coordinador_id=coordinador.uid,
        estado="pendiente",
        comentarios=body.get("comentarios", ""),
    )
    db.add(revision)

    # Cambiar estado del documento
    doc.estado_version = "en_revision"
    db.commit()

    # Notificar al coordinador
    try:
        crear_notificacion(
            db=db,
            usuario_id=coordinador.uid,
            tipo_evento="revision_pendiente",
            descripcion=f"{user_nombre} envió '{doc.nombre}' v{doc.version_actual} para revisión",
            link=f"/r/coordinador/documentos/index.html?doc={doc_id}",
        )
    except Exception as e:
        print(f"Error notificando: {e}")

    return {
        "message": "Documento enviado a revisión",
        "revision_id": revision.id,
        "estado": revision.estado,
    }


@router.post("/revision/{doc_id}/aprobar")
async def aprobar_revision(
    doc_id: int,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    """Coordinador aprueba y envía al Cliente"""
    try:
        body = await request.json()
    except:
        body = {}

    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_email = token_data.get("sub", "")
    user_nombre = token_data.get("nombre", "Usuario")
    roles = token_data.get("roles", [])

    if "coordinador" not in roles and "superadmin" not in roles:
        raise HTTPException(
            status_code=403, detail="Solo el Coordinador puede aprobar revisiones"
        )

    cliente_email = body.get("cliente_email", "")
    if not cliente_email:
        raise HTTPException(status_code=422, detail="email del cliente es requerido")

    # Buscar cliente
    cliente = db.query(User).filter(User.email == cliente_email).first()
    if not cliente:
        raise HTTPException(status_code=404, detail="Cliente no encontrado")

    # Actualizar la revisión pendiente más reciente
    revision = (
        db.query(DocumentoRevision)
        .filter(
            DocumentoRevision.documento_id == doc_id,
            DocumentoRevision.estado == "pendiente",
        )
        .order_by(DocumentoRevision.id.desc())
        .first()
    )

    if revision:
        revision.estado = "aprobado"
        revision.cliente_id = cliente.uid
        revision.comentarios = (
            revision.comentarios or ""
        ) + f"\n[Aprobado por coordinador] {body.get('comentarios', '')}"
        revision.fecha_respuesta = datetime.now()

    doc.estado_version = "aprobada"
    db.commit()

    # Notificar al autor
    try:
        crear_notificacion(
            db=db,
            usuario_id=revision.autor_id if revision else doc.usuario_id,
            tipo_evento="revision_aprobada",
            descripcion=f"'{doc.nombre}' v{doc.version_actual} fue aprobado y enviado al cliente {cliente_email}",
            link=f"/r/autor/documentos/index.html?doc={doc_id}",
        )
    except Exception as e:
        print(f"Error notificando: {e}")

    # Notificar al cliente
    try:
        crear_notificacion(
            db=db,
            usuario_id=cliente.uid,
            tipo_evento="documento_recibido",
            descripcion=f"Has recibido el documento '{doc.nombre}' v{doc.version_actual} (aprobado por coordinador)",
            link=f"/r/cliente/documentos/index.html?doc={doc_id}",
        )
    except Exception as e:
        print(f"Error notificando cliente: {e}")

    return {
        "message": "Documento aprobado y enviado al cliente",
        "estado": doc.estado_version,
    }


@router.post("/revision/{doc_id}/devolver")
async def devolver_revision(
    doc_id: int,
    request: Request,
    db: Session = Depends(get_db),
    token_data: dict = Depends(verify_token),
):
    """Coordinador devuelve el documento al Autor para correcciones"""
    try:
        body = await request.json()
    except:
        body = {}

    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_nombre = token_data.get("nombre", "Usuario")
    roles = token_data.get("roles", [])

    if "coordinador" not in roles and "superadmin" not in roles:
        raise HTTPException(
            status_code=403, detail="Solo el Coordinador puede devolver documentos"
        )

    motivo = body.get("comentarios", "Sin comentarios específicos")

    # Actualizar revisión pendiente
    revision = (
        db.query(DocumentoRevision)
        .filter(
            DocumentoRevision.documento_id == doc_id,
            DocumentoRevision.estado == "pendiente",
        )
        .order_by(DocumentoRevision.id.desc())
        .first()
    )

    if revision:
        revision.estado = "devuelto"
        revision.comentarios = (
            revision.comentarios or ""
        ) + f"\n[Devuelto por coordinador] {motivo}"
        revision.fecha_respuesta = datetime.now()

    # Volver a en_progreso para que el autor pueda editar
    doc.estado_version = "en_progreso"
    db.commit()

    # Notificar al autor
    try:
        crear_notificacion(
            db=db,
            usuario_id=revision.autor_id if revision else doc.usuario_id,
            tipo_evento="revision_devuelta",
            descripcion=f"'{doc.nombre}' v{doc.version_actual} fue devuelto por el coordinador. Motivo: {motivo}",
            link=f"/r/autor/documentos/index.html?doc={doc_id}",
        )
    except Exception as e:
        print(f"Error notificando: {e}")

    return {
        "message": "Documento devuelto al autor",
        "motivo": motivo,
        "estado": doc.estado_version,
    }


@router.get("/revision/pendientes")
async def revisiones_pendientes(
    db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    """Lista documentos pendientes de revisión para el coordinador actual"""
    user_uid = str(token_data.get("user_id") or token_data.get("sub"))
    user_email = token_data.get("sub", "")
    roles = token_data.get("roles", [])

    if "coordinador" not in roles and "superadmin" not in roles:
        user_obj = db.query(User).filter(User.uid == user_uid).first()
        if user_obj:
            db_roles = [r.name for r in user_obj.roles]
            if "coordinador" in db_roles or "superadmin" in db_roles:
                roles = db_roles
        if "coordinador" not in roles and "superadmin" not in roles:
            raise HTTPException(
                status_code=403,
                detail="Solo el Coordinador puede ver revisiones pendientes",
            )

    revisiones = (
        db.query(DocumentoRevision)
        .filter(
            DocumentoRevision.estado == "pendiente",
            DocumentoRevision.coordinador_id == user_uid,
        )
        .order_by(DocumentoRevision.fecha_envio.desc())
        .all()
    )

    result = []
    for rev in revisiones:
        doc = (
            db.query(DocumentoOficina)
            .filter(DocumentoOficina.id == rev.documento_id)
            .first()
        )
        if not doc:
            continue

        autor = db.query(User).filter(User.uid == rev.autor_id).first()

        result.append(
            {
                "revision_id": rev.id,
                "documento_id": rev.documento_id,
                "documento_nombre": doc.nombre,
                "documento_tipo": doc.tipo,
                "version": rev.version,
                "estado": rev.estado,
                "autor_nombre": autor.nombre if autor else rev.autor_id,
                "autor_email": autor.email if autor else "",
                "comentarios": rev.comentarios,
                "fecha_envio": rev.fecha_envio.isoformat() if rev.fecha_envio else None,
            }
        )

    return result


@router.get("/trazabilidad/{doc_id}")
async def trazabilidad_documento(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    """Devuelve la línea de tiempo completa del documento (versiones + revisiones + comparticiones)"""
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    eventos = []

    # Evento: creación
    eventos.append(
        {
            "tipo": "creacion",
            "fecha": doc.creado.isoformat() if doc.creado else None,
            "descripcion": f"Documento creado por usuario {doc.usuario_id}",
            "usuario_id": doc.usuario_id,
        }
    )

    # Eventos: versiones archivadas
    for ver in doc.versiones:
        label = "archivada" if ver.estado != "cerrada" else "cerrada (final)"
        eventos.append(
            {
                "tipo": "version_" + ver.estado,
                "fecha": ver.fecha_creacion.isoformat() if ver.fecha_creacion else None,
                "descripcion": f"Versión v{ver.version} archivada por {ver.creado_por}",
                "version": ver.version,
                "usuario_id": ver.creado_por,
            }
        )
        if ver.fecha_cierre:
            eventos.append(
                {
                    "tipo": "cierre_version",
                    "fecha": ver.fecha_cierre.isoformat(),
                    "descripcion": f"Versión v{ver.version} cerrada",
                    "version": ver.version,
                    "usuario_id": ver.creado_por,
                }
            )

    # Eventos: revisiones (flujo coordinador)
    for rev in doc.revisiones:
        eventos.append(
            {
                "tipo": "enviado_revision",
                "fecha": rev.fecha_envio.isoformat() if rev.fecha_envio else None,
                "descripcion": f"Enviado a revisión (coordinador)",
                "version": rev.version,
                "usuario_id": rev.autor_id,
            }
        )
        if rev.coordinador_visto:
            eventos.append(
                {
                    "tipo": "visto_coordinador",
                    "fecha": rev.coordinador_visto.isoformat()
                    if rev.coordinador_visto
                    else None,
                    "descripcion": "Coordinador vio el documento",
                    "version": rev.version,
                }
            )
        if rev.estado == "aprobado" or rev.estado == "enviado_cliente":
            eventos.append(
                {
                    "tipo": "aprobado",
                    "fecha": rev.fecha_respuesta.isoformat()
                    if rev.fecha_respuesta
                    else None,
                    "descripcion": f"Aprobado por coordinador y enviado al cliente"
                    + (f" (ID: {rev.cliente_id})" if rev.cliente_id else ""),
                    "version": rev.version,
                }
            )
            if rev.cliente_visto:
                eventos.append(
                    {
                        "tipo": "visto_cliente",
                        "fecha": rev.cliente_visto.isoformat()
                        if rev.cliente_visto
                        else None,
                        "descripcion": "Cliente vio el documento",
                        "version": rev.version,
                    }
                )
        elif rev.estado == "devuelto":
            eventos.append(
                {
                    "tipo": "devuelto",
                    "fecha": rev.fecha_respuesta.isoformat()
                    if rev.fecha_respuesta
                    else None,
                    "descripcion": "Devuelto al autor por el coordinador",
                    "version": rev.version,
                }
            )

    # Eventos: comparticiones
    for comp in doc.compartidos:
        eventos.append(
            {
                "tipo": "compartido",
                "fecha": comp.fecha_compartido.isoformat()
                if comp.fecha_compartido
                else None,
                "descripcion": f"Compartido con {comp.email_destino} (permiso: {comp.permiso})",
            }
        )

    # Ordenar por fecha
    eventos.sort(key=lambda e: e.get("fecha") or "")

    return {
        "documento_id": doc.id,
        "documento_nombre": f"{doc.nombre}.{doc.tipo}",
        "version_actual": doc.version_actual,
        "estado_actual": doc.estado_version,
        "eventos": eventos,
    }


# ==================== AUDITORÍA IA ====================


def _extraer_texto_documento(filepath: str, tipo: str) -> str:
    """Extrae texto de un documento Office usando las librerías disponibles"""
    if not os.path.exists(filepath):
        return ""

    try:
        if tipo == "docx":
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif tipo == "xlsx":
            wb = load_workbook(filepath)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        texts.append(" | ".join(cells))
            return "\n".join(texts)
        elif tipo == "pptx":
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
    except Exception as e:
        return f"[Error extrayendo texto: {e}]"

    return ""


@router.post("/auditar/{doc_id}")
async def auditar_documento_ia(
    doc_id: int, db: Session = Depends(get_db), token_data: dict = Depends(verify_token)
):
    """Audita un documento usando IA: extrae texto, compara con cambios solicitados y genera reporte"""
    doc = db.query(DocumentoOficina).filter(DocumentoOficina.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Documento no encontrado")

    # 1. Extraer texto del documento
    filepath = os.path.join(DOCUMENTOS_DIR, doc.filename)
    texto = _extraer_texto_documento(filepath, doc.tipo)

    if not texto or texto.startswith("[Error"):
        raise HTTPException(
            status_code=400, detail="No se pudo extraer texto del documento"
        )

    # 2. Obtener cambios solicitados (desde revisiones/comentarios)
    revisiones = (
        db.query(DocumentoRevision)
        .filter(DocumentoRevision.documento_id == doc_id)
        .order_by(DocumentoRevision.id.desc())
        .all()
    )

    cambios_solicitados = []
    for rev in revisiones:
        if rev.comentarios:
            cambios_solicitados.append(f"[Revisión v{rev.version}] {rev.comentarios}")

    contexto_cambios = (
        "\n".join(cambios_solicitados)
        if cambios_solicitados
        else "No hay cambios solicitados registrados."
    )

    # 3. Llamar a Gemini para la auditoría
    try:
        from app.services.ai.gemini_pool import get_gemini_client

        client = get_gemini_client()
        if not client:
            raise HTTPException(
                status_code=500, detail="No hay API keys de Gemini disponibles"
            )

        prompt = f"""Eres un auditor de documentos educativos. Tu tarea es revisar el siguiente documento y:

1. **RESUMEN**: Haz un resumen breve del contenido del documento (3-5 líneas).
2. **ANÁLISIS DE CAMBIOS**: Compara el contenido con los cambios solicitados abajo. ¿Se implementaron correctamente? ¿Falta algo?
3. **ERRORES Y MEJORAS**: Señala errores gramaticales, ortográficos, de formato, o de contenido.
4. **CALIFICACIÓN**: Asigna una calificación del 1 al 10.
5. **RECOMENDACIONES**: Da recomendaciones concretas para mejorar.

Responde EXACTAMENTE en este formato JSON (sin markdown, sin comillas triples):
{{
  "resumen": "...",
  "analisis_cambios": "...",
  "errores_encontrados": ["...", "..."],
  "calificacion": 7,
  "recomendaciones": ["...", "..."]
}}

**DOCUMENTO:**
{texto[:30000]}

**CAMBIOS SOLICITADOS:**
{contexto_cambios[:5000]}"""

        response = client.models.generate_content(
            model="gemini-2.5-flash", contents=prompt
        )

        resultado_texto = response.text.strip()
        # Limpiar posibles marcadores markdown
        resultado_texto = (
            resultado_texto.replace("```json", "").replace("```", "").strip()
        )

        import json as _json

        try:
            resultado = _json.loads(resultado_texto)
        except:
            resultado = {
                "error": "La IA no devolvió JSON válido",
                "raw": resultado_texto,
            }

        return {
            "documento_id": doc.id,
            "documento_nombre": f"{doc.nombre}.{doc.tipo}",
            "version": doc.version_actual,
            "resultado": resultado,
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error en auditoría IA: {str(e)}")
